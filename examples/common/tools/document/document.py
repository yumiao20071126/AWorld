# coding: utf-8
# Copyright (c) 2025 inclusionAI.

import json
import os
import base64
import tempfile
import subprocess
from typing import Any, Dict, Tuple
from urllib.parse import urlparse

from pydantic import BaseModel

from aworld.config import ToolConfig
from examples.common.tools.tool_action import DocumentExecuteAction
from aworld.core.common import Observation, ActionModel, ActionResult
from aworld.core.tool.base import ToolFactory, Tool
from aworld.logs.util import logger
from examples.common.tools.document.utils import encode_image_from_file, encode_image_from_url
from aworld.utils import import_package, import_packages
from aworld.tools.utils import build_observation


class InputDocument(BaseModel):
    document_path: str | None = None


@ToolFactory.register(name="document_analysis",
                      desc="document analysis",
                      supported_action=DocumentExecuteAction,
                      conf_file_name=f'document_analysis_tool.yaml')
class DocumentTool(Tool):
    def __init__(self, conf: ToolConfig, **kwargs) -> None:
        """Init document tool."""
        import_package('cv2', install_name='opencv-python')
        import_packages(['xmltodict', 'pandas', 'docx2markdown', 'PyPDF2', 'numpy'])
        super(DocumentTool, self).__init__(conf, **kwargs)
        self.cur_observation = None
        self.content = None
        self.keyframes = []
        self.init()
        self.step_finished = True

    def reset(self, *, seed: int | None = None, options: Dict[str, str] | None = None) -> Tuple[
        Observation, dict[str, Any]]:
        super().reset(seed=seed, options=options)

        self.close()
        self.step_finished = True
        return build_observation(observer=self.name(),
                                 ability=DocumentExecuteAction.DOCUMENT_ANALYSIS.value.name), {}

    def init(self) -> None:
        self.initialized = True

    def close(self) -> None:
        pass

    def finished(self) -> bool:
        return self.step_finished

    def do_step(self, actions: list[ActionModel], **kwargs) -> Tuple[Observation, float, bool, bool, Dict[str, Any]]:
        self.step_finished = False
        reward = 0.
        fail_error = ""
        observation = build_observation(observer=self.name(),
                                        ability=DocumentExecuteAction.DOCUMENT_ANALYSIS.value.name)
        info = {}
        try:
            if not actions:
                raise ValueError("actions is empty")
            action = actions[0]
            document_path = action.params.get("document_path", "")
            if not document_path:
                raise ValueError("document path invalid")
            output, keyframes, error = self.document_analysis(document_path)
            observation.content = output
            observation.action_result.append(
                ActionResult(is_done=True,
                             success=False if error else True,
                             content=f"{output}",
                             error=f"{error}",
                             keep=False))
            info['key_frame'] = f"{keyframes}"
            reward = 1.
        except Exception as e:
            fail_error = str(e)
        finally:
            self.step_finished = True
        info["exception"] = fail_error
        info.update(kwargs)
        return (observation, reward, kwargs.get("terminated", False),
                kwargs.get("truncated", False), info)

    def document_analysis(self, document_path):
        import xmltodict
        error = None
        # Initialize content to empty list to avoid None return
        self.content = []
        try:
            if any(document_path.endswith(ext) for ext in [".jpg", ".jpeg", ".png"]):
                parsed_url = urlparse(document_path)
                is_url = all([parsed_url.scheme, parsed_url.netloc])
                if not is_url:
                    base64_image = encode_image_from_file(document_path)
                else:
                    base64_image = encode_image_from_url(document_path)
                self.content = f"data:image/jpeg;base64,{base64_image}"

            if any(document_path.endswith(ext) for ext in ["xls", "xlsx"]):
                try:
                    try:
                        import pandas as pd
                    except ImportError:
                        error = "pandas library not found. Please install pandas: pip install pandas"
                        return self.content, self.keyframes, error

                    excel_data = {}

                    with pd.ExcelFile(document_path) as xls:
                        sheet_names = xls.sheet_names
                        for sheet_name in sheet_names:
                            df = pd.read_excel(xls, sheet_name=sheet_name)
                            sheet_data = df.to_dict(orient='records')
                            excel_data[sheet_name] = sheet_data

                    self.content = json.dumps(excel_data, ensure_ascii=False)
                    logger.info(f"Successfully processed Excel file: {document_path}")
                    logger.info(f"Found {len(sheet_names)} sheets: {', '.join(sheet_names)}")

                except Exception as excel_error:
                    error = str(excel_error)

            if any(document_path.endswith(ext) for ext in ["json", "jsonl", "jsonld"]):
                with open(document_path, "r", encoding="utf-8") as f:
                    self.content = json.load(f)
                f.close()

            if any(document_path.endswith(ext) for ext in ["xml"]):
                data = None
                with open(document_path, "r", encoding="utf-8") as f:
                    data = f.read()
                f.close()

                try:
                    self.content = xmltodict.parse(data)
                    logger.info(f"The extracted xml data is: {self.content}")

                except Exception as e:
                    logger.info(f"The raw xml data is: {data}")
                    error = str(e)
                    self.content = data

            if any(document_path.endswith(ext) for ext in ["doc", "docx"]):
                from docx2markdown._docx_to_markdown import docx_to_markdown
                file_name = os.path.basename(document_path)
                md_file_path = f"{file_name}.md"
                docx_to_markdown(document_path, md_file_path)
                with open(md_file_path, "r") as f:
                    self.content = f.read()
                f.close()

            if any(document_path.endswith(ext) for ext in ["pdf"]):
                # try using pypdf to extract text from pdf
                try:
                    from PyPDF2 import PdfReader

                    # Open file in binary mode for PdfReader
                    f = open(document_path, "rb")
                    reader = PdfReader(f)
                    extracted_text = ""
                    for page in reader.pages:
                        extracted_text += page.extract_text()
                    self.content = extracted_text
                    f.close()
                except Exception as pdf_error:
                    error = str(pdf_error)

            # audio
            if any(document_path.endswith(ext.lower()) for ext in [".mp3", ".wav", ".wave"]):
                try:
                    # audio-> base64
                    with open(document_path, "rb") as audio_file:
                        audio_bytes = audio_file.read()
                        audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                    # ext
                    ext = os.path.splitext(document_path)[1].lower()
                    mime_type = "audio/mpeg" if ext == ".mp3" else "audio/wav"

                    # data URI
                    self.content = f"data:{mime_type};base64,{audio_base64}"
                except Exception as audio_error:
                    error = str(audio_error)
                    logger.error(f"Error processing audio file: {error}")

            # video
            if any(document_path.endswith(ext.lower()) for ext in [".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv"]):
                try:
                    try:
                        import cv2
                        import numpy as np
                    except ImportError:
                        error = "Required libraries not found. Please install opencv-python: pip install opencv-python"
                        return None, None, error

                    # create temp dir
                    temp_dir = tempfile.mkdtemp()

                    # 1.get audio -> base64
                    audio_path = os.path.join(temp_dir, "extracted_audio.mp3")

                    #  get audio by ffmpeg
                    try:
                        subprocess.run([
                            "ffmpeg", "-i", document_path, "-q:a", "0",
                            "-map", "a", audio_path, "-y"
                        ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

                        # audio->base64
                        with open(audio_path, "rb") as audio_file:
                            audio_bytes = audio_file.read()
                            audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                        audio_data_uri = f"data:audio/mpeg;base64,{audio_base64}"
                    except (subprocess.SubprocessError, FileNotFoundError) as e:
                        logger.warning(f"Failed to extract audio: {str(e)}")
                        audio_data_uri = None

                    # 2. get keyframes
                    cap = cv2.VideoCapture(document_path)

                    if not cap.isOpened():
                        raise ValueError(f"Could not open video file: {document_path}")

                    # get video message
                    fps = cap.get(cv2.CAP_PROP_FPS)
                    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                    duration = frame_count / fps if fps > 0 else 0

                    # keyframes policy- per duration/10sï¼Œmax 10
                    keyframes_count = min(10, int(frame_count))
                    frames_interval = max(1, int(frame_count / keyframes_count))

                    self.keyframes = []
                    frame_index = 0

                    while True:
                        ret, frame = cap.read()
                        if not ret:
                            break

                        # per frames_interval save
                        if frame_index % frames_interval == 0:
                            # save JPEG -> base64
                            _, buffer = cv2.imencode(".jpg", frame)
                            img_base64 = base64.b64encode(buffer).decode('utf-8')
                            time_position = frame_index / fps if fps > 0 else 0

                            self.keyframes.append(f"data:image/jpeg;base64,{img_base64}")

                            if len(self.keyframes) >= keyframes_count:
                                break

                        frame_index += 1

                    cap.release()

                    self.content = audio_data_uri
                    logger.info(f"Successfully processed video file: {document_path}")
                    logger.info(f"Extracted {len(self.keyframes)} keyframes and audio track")
                    # clean tmp files
                    try:
                        os.remove(audio_path)
                        os.rmdir(temp_dir)
                    except Exception as cleanup_error:
                        logger.warning(f"Error cleaning up temp files: {str(cleanup_error)}")

                except Exception as video_error:
                    error = str(video_error)
                    logger.error(f"Error processing video file: {error}")

            if any(document_path.endswith(ext) for ext in ["pptx"]):
                try:
                    # Initialize content list and empty keyframes
                    self.content = []
                    self.keyframes = []

                    # Check if file exists
                    if not os.path.exists(document_path):
                        error = f"File does not exist: {document_path}"
                        return self.content, self.keyframes, error

                    # Check if file is readable
                    if not os.access(document_path, os.R_OK):
                        error = f"File is not readable: {document_path}"
                        return self.content, self.keyframes, error

                    # Check file size
                    try:
                        file_size = os.path.getsize(document_path)

                        if file_size == 0:
                            error = "File is empty"
                            return self.content, self.keyframes, error
                    except Exception as size_error:
                        logger.warning(f"Cannot get file size: {str(size_error)}")

                    try:
                        # Import required libraries
                        from pptx import Presentation
                        from PIL import Image, ImageDraw, ImageFont
                        import io
                    except ImportError as import_error:
                        error = f"Missing required libraries: {str(import_error)}. Please install: pip install python-pptx Pillow"
                        return self.content, self.keyframes, error

                    # Create temporary directory for images
                    try:
                        temp_dir = tempfile.mkdtemp()
                    except Exception as temp_dir_error:
                        error = f"Failed to create temporary directory: {str(temp_dir_error)}"
                        return self.content, self.keyframes, error

                    # Open presentation
                    try:
                        presentation = Presentation(document_path)
                        # Get total slides count
                        total_slides = len(presentation.slides)

                        if total_slides == 0:
                            error = "PPTX file does not contain any slides"
                            return self.content, self.keyframes, error

                        # Process each slide
                        for i, slide in enumerate(presentation.slides):

                            # Generate temporary file path for current slide
                            img_path = os.path.join(temp_dir, f"slide_{i + 1}.jpg")

                            # Get slide dimensions
                            try:
                                slide_width = presentation.slide_width
                                slide_height = presentation.slide_height

                                # PPTX dimensions are in EMU (English Metric Unit)
                                # 1 inch = 914400 EMU, 1 cm = 360000 EMU
                                # Convert to pixels (assuming 96 DPI)
                                slide_width_px = int(slide_width / 914400 * 96 * 10)
                                slide_height_px = int(slide_height / 914400 * 96 * 10)

                                # Ensure dimensions are reasonable positive integers
                                slide_width_px = max(1, min(slide_width_px, 4000))  # Limit max width to 4000px
                                slide_height_px = max(1, min(slide_height_px, 3000))  # Limit max height to 3000px

                            except Exception as size_error:
                                # Use default dimensions
                                slide_width_px = 960  # Default width 960px
                                slide_height_px = 720  # Default height 720px

                            # Create blank image
                            try:
                                # Log operation start

                                # Create blank image
                                try:
                                    slide_img = Image.new('RGB', (slide_width_px, slide_height_px), 'white')
                                    draw = ImageDraw.Draw(slide_img)
                                except Exception as img_create_error:
                                    logger.error(
                                        f"Slide {i + 1} blank image creation failed: {str(img_create_error) or 'Unknown error'}")
                                    raise

                                # Draw slide number
                                try:
                                    font = ImageFont.load_default()
                                    draw.text((20, 20), f"Slide {i + 1}/{total_slides}", fill="black", font=font)
                                except Exception as font_error:
                                    logger.warning(f"Failed to draw slide number: {str(font_error) or 'Unknown error'}")

                                # Record shape count
                                try:
                                    shape_count = len(slide.shapes)
                                except Exception as shape_count_error:
                                    logger.warning(
                                        f"Failed to get slide {i + 1} shape count: {str(shape_count_error) or 'Unknown error'}")
                                    shape_count = 0

                                # Try to render shapes on image
                                shape_success_count = 0
                                shape_fail_count = 0

                                try:
                                    for j, shape in enumerate(slide.shapes):
                                        try:
                                            shape_type = type(shape).__name__

                                            # Process images
                                            if hasattr(shape, 'image') and shape.image:
                                                try:
                                                    # Extract image from shape
                                                    image_stream = io.BytesIO(shape.image.blob)
                                                    img = Image.open(image_stream)

                                                    # Calculate position
                                                    left = shape.left
                                                    top = shape.top

                                                    # Paste image onto slide
                                                    slide_img.paste(img, (left, top))
                                                    shape_success_count += 1
                                                except Exception as img_error:
                                                    logger.warning(
                                                        f"Failed to process image {j + 1} in slide {i + 1}: {str(img_error) or 'Unknown error'}")
                                                    if not str(img_error):
                                                        import traceback
                                                        logger.warning(
                                                            f"Image processing stack: {traceback.format_exc()}")
                                                    shape_fail_count += 1

                                            # Process text
                                            elif hasattr(shape, 'text') and shape.text:
                                                try:
                                                    text = shape.text[:30] + "..." if len(
                                                        shape.text) > 30 else shape.text
                                                    # Simple text rendering
                                                    text_left = shape.left
                                                    text_top = shape.top
                                                    draw.text((text_left, text_top), shape.text, fill="black",
                                                              font=font)
                                                    shape_success_count += 1
                                                except Exception as text_error:
                                                    logger.warning(
                                                        f"Failed to process text {j + 1} in slide {i + 1}: {str(text_error) or 'Unknown error'}")
                                                    if not str(text_error):
                                                        import traceback
                                                        logger.warning(
                                                            f"Text processing stack: {traceback.format_exc()}")
                                                    shape_fail_count += 1
                                            else:
                                                logger.info(
                                                    f"Shape {j + 1} in slide {i + 1} is neither image nor text, skipping")
                                        except Exception as shape_error:
                                            if not str(shape_error):
                                                import traceback
                                                logger.warning(f"Shape processing stack: {traceback.format_exc()}")
                                            shape_fail_count += 1
                                except Exception as shapes_iteration_error:
                                    logger.error(
                                        f"Failed while iterating through shapes in slide {i + 1}: {str(shapes_iteration_error) or 'Unknown error'}")
                                    if not str(shapes_iteration_error):
                                        import traceback
                                        logger.error(f"Shape iteration stack: {traceback.format_exc()}")

                                # Save slide image
                                try:
                                    slide_img.save(img_path, 'JPEG')

                                    # Check if image was saved successfully
                                    if not os.path.exists(img_path):
                                        raise ValueError(f"Saved image file does not exist: {img_path}")

                                    file_size = os.path.getsize(img_path)
                                    if file_size == 0:
                                        raise ValueError(
                                            f"Saved image file is empty: {img_path}, size: {file_size} bytes")

                                    # Convert to base64
                                    try:
                                        base64_image = encode_image_from_file(img_path)
                                        self.content.append(f"data:image/jpeg;base64,{base64_image}")
                                    except Exception as base64_error:
                                        error_msg = str(base64_error) or "Unknown base64 conversion error"
                                        if not str(base64_error):
                                            import traceback
                                            logger.error(f"Base64 conversion stack: {traceback.format_exc()}")
                                        raise ValueError(f"Base64 conversion error: {error_msg}")

                                except Exception as save_error:
                                    error_msg = str(save_error) or "Unknown save error"
                                    logger.error(f"Failed to save slide {i + 1} as image: {error_msg}")
                                    if not str(save_error):
                                        import traceback
                                        logger.error(f"Image save stack: {traceback.format_exc()}")
                                    raise ValueError(f"Image save error: {error_msg}")

                            except Exception as slide_render_error:
                                error_msg = str(slide_render_error) or "Unknown rendering error"
                                logger.error(f"Failed to render slide {i + 1}: {error_msg}")
                                if not str(slide_render_error):
                                    import traceback
                                    logger.error(f"Slide rendering stack: {traceback.format_exc()}")
                                # Continue processing next slide, don't interrupt the entire process
                                continue

                    except Exception as pptx_error:
                        error = f"Failed to process PPTX file: {str(pptx_error)}"
                        import traceback

                    # Clean up temporary files
                    try:
                        for file in os.listdir(temp_dir):
                            try:
                                file_path = os.path.join(temp_dir, file)
                                os.remove(file_path)
                            except Exception as file_error:
                                logger.warning(f"Failed to delete temporary file: {str(file_error)}")
                        os.rmdir(temp_dir)
                    except Exception as cleanup_error:
                        logger.warning(f"Failed to clean up temporary files: {str(cleanup_error)}")

                    if len(self.content) > 0:
                        logger.info(f"Extracted {len(self.content)} slides")
                    else:
                        error = error or "Could not extract any slides from PPTX file"
                        logger.error(error)

                except Exception as outer_error:
                    error = f"Error occurred during PPTX file processing: {str(outer_error)}"
                    import traceback

                return self.content, self.keyframes, error

        finally:
            pass

        return self.content, self.keyframes, error
